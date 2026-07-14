"""Extract slide text from .pptx files."""
from io import BytesIO
from typing import List

from pptx import Presentation

from schemas import Slide


def extract_slides(file_bytes: bytes) -> List[Slide]:
    """Parse a .pptx file and return one Slide per slide, joining all of its
    text-frame text (titles, bullets, text boxes) with newlines."""
    prs = Presentation(BytesIO(file_bytes))
    slides: List[Slide] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text.strip())
        slides.append(Slide(index=i, text="\n".join(parts)))
    return slides
